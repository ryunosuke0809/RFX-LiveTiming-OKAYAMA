#!/usr/bin/env python3
"""管理者ページ操作説明の PowerPoint を生成する。"""

from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).with_name("管理者ページ_操作説明.pptx")

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(0x0C, 0x0C, 0x0F)
BAR = RGBColor(0x18, 0x18, 0x1B)
CARD = RGBColor(0x1C, 0x1C, 0x21)
AMBER = RGBColor(0xD9, 0x77, 0x06)
WHITE = RGBColor(0xFA, 0xFA, 0xF9)
MUTED = RGBColor(0xA1, 0xA1, 0xAA)
DIM = RGBColor(0x71, 0x71, 0x7A)
GREEN = RGBColor(0x34, 0xD3, 0x99)
FONT = "Hiragino Sans"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def set_run_font(run, size_pt: float, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = etree.SubElement(rPr, f"{{{NS_A}}}{tag}")
        el.set("typeface", FONT)


def add_rect(slide, l, t, w, h, fill: RGBColor, line: RGBColor | None = None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Emu(6350)
    sh.shadow.inherit = False
    return sh


def add_round(slide, l, t, w, h, fill: RGBColor):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.adjustments[0] = 0.08
    sh.shadow.inherit = False
    return sh


def tb(slide, l, t, w, h, text: str, size: float, color: RGBColor, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, size, color, bold)
    return box


def bullets(slide, l, t, w, h, items: list[str], size: float = 18):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "・  " + item
        set_run_font(run, size, WHITE, False)
    return box


def add_table(slide, l, t, w, rows: list[list[str]], col_w: list[int], header=True):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, l, t, w, Inches(0.42 * n_rows))
    table = table_shape.table
    total = sum(col_w)
    for i, cw in enumerate(col_w):
        table.columns[i].width = int(w * cw / total)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            is_head = header and r == 0
            set_run_font(run, 13 if is_head else 13, WHITE if is_head else WHITE, is_head)
            cell.text_frame.word_wrap = True
            fill = AMBER if is_head else (CARD if r % 2 == 0 else BAR)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
    return table_shape


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    return slide


def chrome(slide, title: str, page: int, total: int):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.92), BAR)
    add_rect(slide, 0, 0, Inches(0.12), Inches(0.92), AMBER)
    tb(slide, Inches(0.45), Inches(0.22), Inches(11.5), Inches(0.5), title, 24, WHITE, True)
    add_rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), BAR)
    tb(
        slide,
        Inches(0.45),
        SLIDE_H - Inches(0.36),
        Inches(10),
        Inches(0.32),
        "MOLA Timing  ·  岡山国際サーキット  ·  管理者専用  ·  この URL は一般公開しない",
        10,
        DIM,
        False,
    )
    tb(
        slide,
        Inches(11.4),
        SLIDE_H - Inches(0.36),
        Inches(1.6),
        Inches(0.32),
        f"{page}  /  {total}",
        10,
        DIM,
        False,
        PP_ALIGN.RIGHT,
    )


def card(slide, l, t, w, h, title: str, body: str, accent: RGBColor = AMBER):
    add_round(slide, l, t, w, h, CARD)
    add_rect(slide, l, t, Inches(0.1), h, accent)
    tb(slide, l + Inches(0.28), t + Inches(0.18), w - Inches(0.45), Inches(0.36), title, 15, accent, True)
    box = slide.shapes.add_textbox(l + Inches(0.28), t + Inches(0.52), w - Inches(0.45), h - Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    set_run_font(run, 14, WHITE, False)


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = 14

    # 1 title
    s = new_slide(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_rect(s, 0, 0, Inches(0.16), SLIDE_H, AMBER)
    tb(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.4), "OKAYAMA INTERNATIONAL CIRCUIT", 14, AMBER, True)
    tb(s, Inches(0.8), Inches(2.15), Inches(11.5), Inches(1.1), "管理者ページ 操作説明", 40, WHITE, True)
    tb(
        s,
        Inches(0.8),
        Inches(3.4),
        Inches(11.5),
        Inches(0.6),
        "MOLA Timing  ·  Live Timing / 履歴データの公開管理",
        18,
        MUTED,
        False,
    )
    tb(
        s,
        Inches(0.8),
        Inches(5.4),
        Inches(11.5),
        Inches(0.8),
        "2026-08-22  提供版\n対象: 運営・事務局の管理者    ※ この資料と URL は関係者限り",
        14,
        DIM,
        False,
    )

    # 2 overview
    s = new_slide(prs)
    chrome(s, "この画面でできること", 2, total)
    card(
        s,
        Inches(0.45),
        Inches(1.2),
        Inches(4.0),
        Inches(4.6),
        "履歴データ",
        "お客様が見る過去リザルトを、日付またはセッション単位で公開画面から外す／戻す。\n\n大会名・セッション名など、表示名だけ変えることもできます。\n\n非表示にしても計測データは消えません。",
    )
    card(
        s,
        Inches(4.65),
        Inches(1.2),
        Inches(4.0),
        Inches(4.6),
        "Live表示",
        "いま動いている Live Timing 表の列を、並び替え・名称変更・表示／非表示できます。\n\nCar / Behind などのプルダウンの文字や初期表示も変えられます。\n\nタイムの数値そのものは変わりません。",
        GREEN,
    )
    card(
        s,
        Inches(8.85),
        Inches(1.2),
        Inches(4.0),
        Inches(4.6),
        "ユーザー",
        "管理者アカウントの追加、パスワード変更、削除ができます。\n\n管理画面では位置情報の許可は不要です。",
        RGBColor(0x60, 0xA5, 0xFA),
    )

    # 3 URL
    s = new_slide(prs)
    chrome(s, "開くアドレス（一般公開しない）", 3, total)
    tb(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4), "管理者専用", 14, AMBER, True)
    add_round(s, Inches(0.45), Inches(1.7), Inches(12.4), Inches(1.15), CARD)
    tb(
        s,
        Inches(0.7),
        Inches(1.95),
        Inches(12.0),
        Inches(0.7),
        "https://oic-timing-admin.mola-timing-okayama.com/",
        22,
        WHITE,
        True,
    )
    tb(s, Inches(0.5), Inches(3.15), Inches(12.3), Inches(0.35), "お客様・関係者が使う画面（こちらは公開してよい）", 14, MUTED, True)
    add_table(
        s,
        Inches(0.45),
        Inches(3.55),
        Inches(12.4),
        [
            ["URL", "誰向け", "管理画面の変更"],
            ["https://mola-timing-okayama.com/", "一般（場内・位置情報あり）", "Live と履歴にそのまま出る"],
            ["https://oic-private.mola-timing-okayama.com/", "関係者（位置情報なし）", "同じ Live / 同じ履歴"],
        ],
        [5, 4, 3],
    )
    tb(
        s,
        Inches(0.5),
        Inches(6.35),
        Inches(12.3),
        Inches(0.4),
        "SNS・掲示・お客様案内に、管理者用 URL を載せないでください。一般向けからは管理画面は開けません。",
        14,
        MUTED,
    )

    # 4 login
    s = new_slide(prs)
    chrome(s, "ログイン", 4, total)
    bullets(
        s,
        Inches(0.5),
        Inches(1.2),
        Inches(7.2),
        Inches(2.2),
        [
            "上の管理者 URL をブラウザで開く（自動で管理画面へ移動します）",
            "ユーザー名とパスワードを入力し、「ログイン」を押す",
            "使い始めたら、ユーザー画面でパスワードを変更してください",
        ],
        17,
    )
    add_round(s, Inches(8.0), Inches(1.2), Inches(4.85), Inches(2.55), CARD)
    add_rect(s, Inches(8.0), Inches(1.2), Inches(0.1), Inches(2.55), AMBER)
    tb(s, Inches(8.3), Inches(1.4), Inches(4.4), Inches(0.35), "初回アカウント（2026-08-22）", 13, AMBER, True)
    tb(s, Inches(8.3), Inches(1.9), Inches(4.4), Inches(1.5), "ユーザー名\n   timing_oic\n\nパスワード\n   timing_oic", 16, WHITE, False)
    add_table(
        s,
        Inches(0.45),
        Inches(4.0),
        Inches(12.4),
        [
            ["注意", "内容"],
            ["セッション", "操作中はログインが維持されます（目安 12 時間。操作が無いと切れます）"],
            ["切れたとき", "同じ画面で再度ログインしてください"],
            ["渡し方", "この資料ごと広く配らず、アカウントは口頭または別紙が安全です"],
        ],
        [2, 8],
    )

    # 5 nav
    s = new_slide(prs)
    chrome(s, "画面の見方", 5, total)
    tb(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.4), "ログイン後、画面上部のナビで切り替えます。", 16, MUTED)
    items = [
        ("履歴データ", "過去リザルトの公開／非公開と表示名"),
        ("Live表示", "いま動いている Live Timing 表の見た目"),
        ("ユーザー", "管理者アカウントの追加・変更・削除"),
        ("右上の名前", "いまログインしている人"),
        ("ログアウト", "使い終わったら押す（共用 PC では特に）"),
    ]
    for i, (title, body) in enumerate(items):
        y = Inches(1.75) + Inches(0.9) * i
        add_round(s, Inches(0.45), y, Inches(12.4), Inches(0.8), CARD)
        add_rect(s, Inches(0.45), y, Inches(0.1), Inches(0.8), AMBER)
        tb(s, Inches(0.85), y + Inches(0.12), Inches(3.2), Inches(0.55), title, 18, WHITE, True)
        tb(s, Inches(4.2), y + Inches(0.18), Inches(8.3), Inches(0.5), body, 16, MUTED, False)

    # 6 archive intro
    s = new_slide(prs)
    chrome(s, "履歴データ  —  消えるわけではありません", 6, total)
    tb(
        s,
        Inches(0.5),
        Inches(1.25),
        Inches(12.3),
        Inches(0.9),
        "お客様が見る履歴（カレンダー・リザルト）に出す／出さないを決めます。\n「非表示にする」は削除ではありません。データはサーバーに残り、いつでも戻せます。",
        18,
        WHITE,
    )
    add_table(
        s,
        Inches(0.45),
        Inches(2.5),
        Inches(12.4),
        [
            ["やりたいこと", "操作"],
            ["テスト走行を公開したくない", "そのセッションを「非表示にする」"],
            ["その日まるごと出したくない", "日付を「非表示にする」"],
            ["「AM走行」を分かりやすい名前にしたい", "「表示名」でセッション名を変更して保存"],
            ["間違えて隠した", "「再表示する」（データは残っています）"],
        ],
        [4, 6],
    )

    # 7 archive ops
    s = new_slide(prs)
    chrome(s, "履歴データ  —  日付とセッション", 7, total)
    card(
        s,
        Inches(0.45),
        Inches(1.2),
        Inches(6.15),
        Inches(5.4),
        "日付ごと（その日全部）",
        "1. 「履歴データ」を開く\n2. 日付一覧を確認（新しい日付が上）\n3. 出さない日は「非表示にする」\n4. カレンダーや履歴からその日が消える\n5. 戻すときは「再表示する」\n\n日付をタップすると、その日のセッション一覧が開きます。\n日付を非表示にすると、中のセッションもまとめて公開されません。",
    )
    card(
        s,
        Inches(6.8),
        Inches(1.2),
        Inches(6.05),
        Inches(5.4),
        "セッションごと（例: AM走行だけ）",
        "1. 対象の日付を選ぶ\n2. セッション名の右の「非表示にする」\n3. そのセッションだけ履歴から消える\n4. 戻すときは「再表示する」\n\n日付は出したまま、一部セッションだけ隠す、という使い方ができます。",
        GREEN,
    )

    # 8 names
    s = new_slide(prs)
    chrome(s, "履歴データ  —  表示名の変更", 8, total)
    bullets(
        s,
        Inches(0.5),
        Inches(1.2),
        Inches(12.3),
        Inches(1.8),
        [
            "対象セッションの「表示名」を押す",
            "カテゴリ / ラウンド / セッション / 大会 を必要に応じて書き換えて「保存」",
            "空欄にして保存すると、上書きをやめて計測データそのままの名前に戻る",
            "変更は履歴の一覧・詳細・CSV のファイル名にも反映されます",
        ],
        17,
    )
    add_table(
        s,
        Inches(0.45),
        Inches(3.5),
        Inches(12.4),
        [
            ["項目", "例"],
            ["カテゴリ", "クラス・カテゴリ名"],
            ["ラウンド", "第○戦 など"],
            ["セッション", "AM走行、公式予選、決勝 など"],
            ["大会", "イベント名"],
        ],
        [3, 7],
    )

    # 9 live save
    s = new_slide(prs)
    chrome(s, "Live表示  —  保存するまでお客様側は変わりません", 9, total)
    add_round(s, Inches(0.45), Inches(1.25), Inches(12.4), Inches(1.7), CARD)
    add_rect(s, Inches(0.45), Inches(1.25), Inches(0.1), Inches(1.7), AMBER)
    tb(
        s,
        Inches(0.8),
        Inches(1.5),
        Inches(11.8),
        Inches(1.3),
        "並びや名前をいじっただけでは、Live は変わりません。\n右上の「保存して Live に反映」を押した時点で、開いている Live ページにもすぐ届きます（再読み込み不要）。",
        18,
        WHITE,
    )
    bullets(
        s,
        Inches(0.5),
        Inches(3.25),
        Inches(12.3),
        Inches(3.2),
        [
            "未保存のときは「未保存の変更あり」と出ます",
            "保存せず別メニューへ移動すると、変更は破棄されます",
            "元に戻す: 「既定に戻す」→「保存して Live に反映」",
            "確認: 保存したあと、別タブで一般向けまたは関係者向けの Live を見る",
            "届かないときは、その Live ページを再読み込みしてください",
        ],
        17,
    )

    # 10 live ops
    s = new_slide(prs)
    chrome(s, "Live表示  —  並び・名称・表示／非表示", 10, total)
    ops = [
        ("並び替え", "各行の左の ▲ / ▼ で順番を変えます。上のプレビューで Live に出る並びを確認できます。"),
        ("表示／非表示", "「表示」チェックを外すとその列は Live に出ません。例: セクターやピット回数を出さない。"),
        ("名称", "右側の入力欄が表ヘッダーの名前です。例: P → 順位、No. → 車番。24 文字まで。空でも可。"),
        ("注意", "決勝以外でも使う列は、少なくとも 1 つは表示にしてください。「決勝のみ」の列（▲▼）は決勝のときだけ出ます。"),
    ]
    for i, (title, body) in enumerate(ops):
        y = Inches(1.2) + Inches(1.3) * i
        add_round(s, Inches(0.45), y, Inches(12.4), Inches(1.18), CARD)
        add_rect(s, Inches(0.45), y, Inches(0.1), Inches(1.18), AMBER)
        tb(s, Inches(0.85), y + Inches(0.14), Inches(12.0), Inches(0.35), title, 16, AMBER, True)
        tb(s, Inches(0.85), y + Inches(0.5), Inches(12.0), Inches(0.55), body, 15, WHITE, False)

    # 11 dropdown + columns
    s = new_slide(prs)
    chrome(s, "Live表示  —  プルダウンと列の意味", 11, total)
    tb(
        s,
        Inches(0.5),
        Inches(1.15),
        Inches(12.3),
        Inches(0.7),
        "一部の列は、お客様がヘッダーをタップして表示を切り替えられます。「プルダウン」から編集します。出す項目は 1 つ以上残してください。",
        15,
        MUTED,
    )
    add_table(
        s,
        Inches(0.45),
        Inches(1.9),
        Inches(6.1),
        [
            ["操作", "意味"],
            ["出す", "お客様のメニューにその項目を出す"],
            ["名称", "メニューに出る文字（Behind → 先頭差 など）"],
            ["初期", "ページを開いたときに最初に出す項目"],
        ],
        [2, 5],
    )
    add_table(
        s,
        Inches(6.75),
        Inches(1.9),
        Inches(6.1),
        [
            ["キー", "中身"],
            ["pos / pic / nr", "総合順位 / クラス内順位 / 車番"],
            ["driver / car / class", "名前 / 車両・チーム / クラス"],
            ["laps / gap / best", "周回・時刻 / 差 / ベスト"],
            ["s1 s2 s3 / pit / chg", "セクター / ピット / 順位変動（決勝）"],
        ],
        [3, 5],
    )
    tb(
        s,
        Inches(0.5),
        Inches(5.85),
        Inches(12.3),
        Inches(0.8),
        "プルダウンを 1 項目だけにすると、お客様側は切替メニューではなく、その名前の固定見出しになります。",
        15,
        MUTED,
    )

    # 12 users
    s = new_slide(prs)
    chrome(s, "ユーザー（管理者アカウント）", 12, total)
    card(
        s,
        Inches(0.45),
        Inches(1.2),
        Inches(4.0),
        Inches(5.4),
        "追加",
        "「ユーザー」を開き、下の「管理者を追加」に入れて「追加」。\n\nユーザー名: 英数字と . _ - のみ。3〜32 文字。日本語は不可。\n\nパスワード: 10 文字以上。",
    )
    card(
        s,
        Inches(4.65),
        Inches(1.2),
        Inches(4.0),
        Inches(5.4),
        "パスワード変更",
        "対象の「パスワード変更」→ 新しいパスワード →「変更」。\n\nその人のログインは切れます。入り直しが必要です。\n\n自分のパスワードを変えた場合は、この画面からもログアウトします。",
        GREEN,
    )
    card(
        s,
        Inches(8.85),
        Inches(1.2),
        Inches(4.0),
        Inches(5.4),
        "削除",
        "対象の「削除」→ 確認ダイアログで実行。\n\n自分自身は削除できません。\n\n管理者が 1 人だけのときは削除できません。先に別の人を追加してください。",
        RGBColor(0xF8, 0x71, 0x71),
    )

    # 13 trouble
    s = new_slide(prs)
    chrome(s, "困ったとき", 13, total)
    add_table(
        s,
        Inches(0.35),
        Inches(1.15),
        Inches(12.6),
        [
            ["症状", "確認すること"],
            ["ログインできない", "ユーザー名／パスワードの打ち間違い。大文字小文字も一致させる"],
            ["管理 URL が Live になる", "アドレスが oic-timing-admin か。一般向けでは管理画面は出ません"],
            ["履歴から消したのにデータが心配", "非表示は公開から外しているだけ。「再表示する」で戻る"],
            ["Live の列が変わらない", "「保存して Live に反映」を押したか。未保存のまま移動していないか"],
            ["プルダウンが保存できない", "その列を表示しているなら、「出す」を 1 つ以上オンにする"],
            ["決勝の ▲▼ が出ない", "列 chg が非表示でないか。予選・練習では仕様上出ません"],
            ["ログイン画面に戻った", "しばらく操作が無かったか、パスワード変更後。入り直す"],
        ],
        [4, 8],
    )

    # 14 ops
    s = new_slide(prs)
    chrome(s, "運用上のお願い", 14, total)
    rules = [
        "管理 URL とパスワードは、必要な担当者だけが持つ",
        "使い終わったらログアウトする（共用 PC の場合は特に）",
        "履歴の非表示は「消す」ではなく「お客様に見せない」操作として使う",
        "Live の見た目を大きく変えるときは、本番セッションの直前ではなく余裕のある時間に保存して確認する",
        "分からない操作は、試す前にこの資料の該当ページを確認する",
        "問い合わせは、この資料をお渡しした担当者まで",
    ]
    for i, text in enumerate(rules):
        y = Inches(1.25) + Inches(0.85) * i
        add_round(s, Inches(0.45), y, Inches(12.4), Inches(0.75), CARD)
        add_rect(s, Inches(0.45), y, Inches(0.72), Inches(0.75), AMBER)
        tb(s, Inches(0.52), y + Inches(0.18), Inches(0.58), Inches(0.45), f"{i + 1}", 18, WHITE, True, PP_ALIGN.CENTER)
        tb(s, Inches(1.4), y + Inches(0.18), Inches(11.2), Inches(0.45), text, 16, WHITE, False)

    prs.save(OUT)
    print(f"wrote {OUT}")


if __name__ == "__main__":
    build()
